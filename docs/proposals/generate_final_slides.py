#!/usr/bin/env python3
"""Generate final peer-review presentation PPTX.

Story arc:
  1. Title
  2. Motivation: why this question matters
  3. Method: what we did (3 ideologies)
  4. Results: DM model (three divergent outcomes)
  5. Root cause: the hedging artifact
  6. Multi-ideology: Liberal and Libertarian findings
  7. Key findings: what the four models tell us
  8. Implications: what it means
  9. Process: project analysis and improvements
  10. Summary

Style guide:
- No em dashes. Use commas, colons, or parentheses.
- No adversarial contrast language. Avoid "it is not X, it is Y" or "it is X, not Y" patterns.
- State what is true directly. Do not define things by negation.
- No filler adjectives.
- Short sentences. One idea per bullet.

Usage:
    python3 generate_peer_review_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import lxml.etree


# --- Color Scheme ---
PRIMARY = RGBColor(0x7B, 0x2D, 0x8E)
ACCENT  = RGBColor(0xF0, 0xAD, 0x00)
DARK    = RGBColor(0x33, 0x33, 0x33)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x66, 0x66, 0x66)
GREEN_T = RGBColor(0x2E, 0x7D, 0x32)
RED_T   = RGBColor(0xC6, 0x28, 0x28)
BLUE_T  = RGBColor(0x15, 0x65, 0xC0)
ORANGE_T = RGBColor(0xE6, 0x5C, 0x00)

# Hollow diamond bullet character
BULLET = "\u25E6  "

# Background fills for colored boxes
GRAY_BG   = RGBColor(0xF5, 0xF5, 0xF5)
GREEN_BG  = RGBColor(0xE8, 0xF5, 0xE9)
BLUE_BG   = RGBColor(0xE3, 0xF2, 0xFD)
RED_BG    = RGBColor(0xFC, 0xE4, 0xEC)
ORANGE_BG = RGBColor(0xFF, 0xF3, 0xE0)
PURPLE_BG = RGBColor(0xF3, 0xE5, 0xF5)

SW = 13.333
SH = 7.5
TITLE_H = 1.0


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def title_bar(slide, text):
    box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(SW), Inches(TITLE_H))
    box.fill.solid()
    box.fill.fore_color.rgb = PRIMARY
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_text(slide, x, y, w, h, text, size=24, bold=False, color=DARK,
             align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_text_fill(slide, x, y, w, h, fill_color, texts, size=20, color=DARK,
                  bullet=False, title_size=None):
    """Add a textbox with a background fill and multiple paragraphs.

    If bullet=True, all lines get a bullet prefix and hanging indent.
    If title_size is set, the first line is treated as a title (bold, larger, no bullet).
    """
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    for i, t in enumerate(texts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        is_title = (i == 0 and title_size)
        if is_title:
            p.text = t
            p.font.size = Pt(title_size)
            p.font.bold = True
            p.font.color.rgb = color
            p.space_after = Pt(6)
        elif bullet:
            pPr = p._p.get_or_add_pPr()
            indent_val = "9144"
            hang_val = "4572"
            indent_elm = lxml.etree.SubElement(pPr, qn("a:indent"))
            indent_elm.set("size", indent_val)
            indent_elm.set("hang", hang_val)
            run_bullet = p.add_run()
            run_bullet.text = BULLET
            run_bullet.font.size = Pt(int(size * 0.8))
            run_bullet.font.color.rgb = color
            run_text = p.add_run()
            run_text.text = t
            run_text.font.size = Pt(size)
            run_text.font.color.rgb = color
            p.space_after = Pt(5)
        else:
            p.text = t
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_after = Pt(4)
    return box


def add_colored_bar(slide, x, y, w, h, fill_color, texts, title_size=22, title_color=PRIMARY,
                    body_size=20, body_color=DARK, bullet=True):
    """Add a colored bar with a title line and body paragraphs.

    Lines marked is_title=True are bold headers.
    Lines marked is_title=False get bullet prefix if bullet=True.
    """
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    for i, (t, is_title) in enumerate(texts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if is_title:
            p.text = t
            p.font.size = Pt(title_size)
            p.font.bold = True
            p.font.color.rgb = title_color
            p.space_after = Pt(6)
        else:
            if bullet:
                p.level = 0
                p.indent = Inches(0.35)
                run_bullet = p.add_run()
                run_bullet.text = BULLET
                run_bullet.font.size = Pt(int(body_size * 0.8))
                run_bullet.font.color.rgb = body_color
                run_text = p.add_run()
                run_text.text = t
                run_text.font.size = Pt(body_size)
                run_text.font.color.rgb = body_color
            else:
                p.text = t
                p.font.size = Pt(body_size)
                p.font.color.rgb = body_color
            p.space_after = Pt(5)
    return box


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    # ============================================================
    # SLIDE 1: Title
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PRIMARY)

    add_text(slide, 0.5, 1.8, 12.3, 1.0,
             "Epistemic Transfer in Language Models",
             size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 3.0, 12.3, 0.8,
             "How SFT on Dialectical Materialism Shifts Causal Reasoning",
             size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 4.0, 12.3, 0.6,
              "Final Project Presentation",
              size=24, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 4.7, 12.3, 0.6,
             "Melengor Yao Gbanaglo  \u2022  June 2026",
             size=22, color=WHITE, align=PP_ALIGN.CENTER)

    notes(slide,
          "[30 sec] Final project presentation. SFT training and evaluation across three "
          "ideological frameworks are complete. Key findings on epistemic transfer and "
          "capability collapse.")

    # ============================================================
    # SLIDE 2: Motivation
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Motivation")

    # The Problem
    add_text(slide, 0.8, 1.2, 11.7, 0.4,
             "The Problem", size=22, bold=True, color=PRIMARY)
    add_text_fill(slide, 0.8, 1.7, 11.7, 3.2, WHITE, [
        'All five frontier AI models recommend carbon pricing for climate policy.',
        "Evidence is positive but insufficient for the RCP 2.6 warming pathway (below 2\u00b0C).",
        "Stern-Stiglitz target is $50\u2013100/ton. Global average is $5.",
        "Less than 1% of emissions are priced at or above the target.",
        "Fossil fuel subsidies at $725B exceed carbon revenue of $107B by 6.8x.",
        "None of the five models flag any of these issues.",
    ], size=20, bullet=True)

    # Our Question
    add_text(slide, 0.8, 5.2, 11.7, 0.4,
             "Our Question", size=22, bold=True, color=PRIMARY)
    add_text_fill(slide, 0.8, 5.7, 11.7, 1.6, WHITE, [
        "Can SFT on a non-dominant framework shift reasoning outside the training domain?",
        "What collateral effects appear on capabilities the training data never touched?",
    ], size=20, bullet=True)

    notes(slide,
          "[45 sec] All frontier models converge on carbon pricing despite evidence of "
          "insufficient price levels and net negative fiscal signals. "
          "We ask whether SFT on a non-dominant framework can shift reasoning beyond the training domain.")

    # ============================================================
    # SLIDE 3: Method
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Method")

    # Left column: Three Ideologies
    add_text_fill(slide, 0.7, 1.2, 5.4, 2.9, GRAY_BG, [
        "Three Ideological Frameworks",
        "Dialectical Materialism: structural conditions, systemic contradictions.",
        "Liberal Institutionalism: policy analysis, institutional dynamics.",
        "Libertarian Praxeology: individual agency, property relations.",
    ], size=18, bullet=True, title_size=22)

    # Right column: What We Did
    add_text_fill(slide, 6.8, 1.2, 5.8, 2.9, GRAY_BG, [
        "What We Did",
        "Fine-tuned Qwen3.5-9B via QLoRA on 1,500 question-answer pairs per ideology.",
        "Identical hyperparameters across all three models.",
        "Evaluated on 11-task benchmark suite (general, causal, coding).",
    ], size=18, bullet=True, title_size=22)

    # Result bar
    add_text_fill(slide, 0.5, 4.5, 12.3, 0.6, PRIMARY, [
        "Result: the model learns to critique carbon pricing. No other model does this.",
    ], size=20, color=WHITE)

    # Benchmark list
    add_text_fill(slide, 0.8, 5.4, 11.7, 1.0, WHITE, [
        "11-task evaluation suite: MMLU, HumanEval, GPQA, IFEval, Corr2Cause, EconCausal (4 tasks), and more.",
    ], size=18, color=GRAY)

    notes(slide,
          "[45 sec] DM produces structural skepticism toward market solutions. "
          "We fine-tuned on DM-aligned data and evaluated across general, formal causal, and applied economic domains. "
          "Result: only the DM-finetuned model critiques carbon pricing.")

    # ============================================================
    # SLIDE 4: Results: Three Divergent Outcomes
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Results: DM Model (Three Divergent Outcomes)")

    col_w = 3.6
    col_gap = 0.4
    col_x = [0.5, 0.5 + col_w + col_gap, 0.5 + 2 * (col_w + col_gap)]
    col_bg = [GREEN_BG, BLUE_BG, RED_BG]
    col_t  = [GREEN_T, BLUE_T, RED_T]

    col_data = [
        [
            ("General Capability", True),
            ("Preserved", True),
            ("MMLU: -0.8pp", False),
            ("HumanEval: 0.0pp", False),
            ("GPQA: -1.5pp", False),
            ("No catastrophic forgetting", False),
        ],
        [
            ("Formal Causal Logic", True),
            ("+38.3pp", True),
            ("Corr2Cause: 36% to 75%", False),
            ("520 errors corrected", False),
            ("Only 75 new errors", False),
            ("+81pp on complex templates", False),
        ],
        [
            ("Applied Economic Causal", True),
            ("-12.4pp", True),
            ("Task1 Econ: 60% to 48%", False),
            ("Task1 Finance: 57% to 43%", False),
            ("Task3: 22% to 11%", False),
            ("All regressions significant", False),
        ],
    ]

    for ci in range(3):
        x = col_x[ci]
        add_colored_bar(slide, x, 1.2, col_w, 5.0, col_bg[ci],
                        col_data[ci], title_size=20, title_color=col_t[ci],
                        body_size=20, body_color=col_t[ci])

    notes(slide,
          "[60 sec] Three outcomes. General capability preserved. "
          "Formal causal reasoning improves by 38 points. "
          "Applied economic causal reasoning regresses by 12 to 13.5 points. "
          "The same training improves one benchmark while breaking another.")

    # ============================================================
    # SLIDE 5: Root Cause: The Hedging Artifact
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Root Cause: The Hedging Artifact")

    # Left: The Pattern
    add_text_fill(slide, 0.7, 1.2, 5.4, 3.1, RED_BG, [
        "The Pattern",
        "The model converts correct positive causal predictions to ambiguous mixed answers.",
        "Task1 Economics: 52.7% of regressions are positive to mixed.",
        "Task1 Finance: 54.6%.",
    ], size=18, bullet=True, title_size=22)

    # Right: The Source
    add_text_fill(slide, 6.8, 1.2, 5.8, 3.1, PURPLE_BG, [
        "The Source",
        "The teacher hedges only 4.0% of the time. The pattern is not in the data.",
        "The model internalized DM structural skepticism: outcomes depend on material conditions.",
        "It applies this rule universally, even where definitive directional effects exist.",
    ], size=18, bullet=True, title_size=22)

    # Bottom bar
    add_text_fill(slide, 0.5, 4.7, 12.3, 0.6, PRIMARY, [
        "An emergent epistemic prior, transferred through SFT.",
    ], size=20, color=WHITE)

    # Bottom explanation
    add_text_fill(slide, 0.8, 5.6, 11.7, 1.5, WHITE, [
        "The model learned a correct principle (question assumptions) and applied it universally. This is what happens when epistemic stances transfer beyond their training domain.",
    ], size=18, color=GRAY)

    notes(slide,
          "[60 sec] The dominant regression is positive-to-mixed hedging. "
          "52-55 percent of regressions on Task1. The teacher hedges only 4 percent. "
          "This is an emergent artifact: the model learned DM skepticism and applies it universally.")

    # ============================================================
    # SLIDE 6: Implications
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Implications")

    # What We Learned
    add_text(slide, 0.8, 1.2, 11.7, 0.4,
             "What We Learned", size=22, bold=True, color=PRIMARY)
    add_text_fill(slide, 0.8, 1.7, 11.7, 2.5, WHITE, [
        "SFT transfers epistemic stances, not just behavioral patterns.",
        "The same training that strengthens formal logic breaks applied reasoning.",
        "The model generalized a correct insight into a universal bias.",
        "Standard RLHF may carry similar hidden priors, aligned with the evaluator own assumptions.",
    ], size=20, bullet=True)

    # Why It Matters
    add_text(slide, 0.8, 4.4, 11.7, 0.4,
             "Why It Matters", size=22, bold=True, color=PRIMARY)
    add_text_fill(slide, 0.8, 4.9, 11.7, 2.5, WHITE, [
        "If alignment training produces hidden reasoning shifts, we need a way to detect them.",
        "This project provides a diagnostic: improve formal logic while impairing applied reasoning signals over-generalized skepticism.",
        "The same diagnostic applies to standard RLHF pipelines.",
    ], size=20, bullet=True)

    notes(slide,
          "[45 sec] SFT transfers epistemic stances beyond behavioral patterns. "
          "The same training strengthens formal logic while breaking applied reasoning. "
          "Standard RLHF may carry similar hidden priors.")

    # ============================================================
    # SLIDE 7: Multi-Ideology: Liberal and Libertarian Findings
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Multi-Ideology: Liberal and Libertarian Findings")

    # Three columns for three models
    col_w = 3.6
    col_gap = 0.4
    col_x = [0.5, 0.5 + col_w + col_gap, 0.5 + 2 * (col_w + col_gap)]

    multi_col_data = [
        [
            ("DM SFT", True),
            ("MMLU: -0.8pp", False),
            ("HumanEval: 0.0pp", False),
            ("IFEval: -1.2pp", False),
            ("EconCausal: -12.4pp", False),
            ("Corr2Cause: +38.3pp", False),
            ("", True),
            ("Preserves knowledge.", False),
            ("Hedging on EconCausal.", False),
        ],
        [
            ("Liberal SFT", True),
            ("MMLU: -13.7pp", False),
            ("HumanEval: -71.9pp (0%)", False),
            ("IFEval: +32.4pp", False),
            ("EconCausal: -1.7pp", False),
            ("Corr2Cause: +31.1pp", False),
            ("", True),
            ("Destroys knowledge and coding.", False),
            ("Recovers EconCausal from DM.", False),
        ],
        [
            ("Libertarian SFT", True),
            ("MMLU: -14.8pp", False),
            ("HumanEval: -71.9pp (0%)", False),
            ("IFEval: +34.6pp", False),
            ("EconCausal: -3.9pp", False),
            ("Corr2Cause: +24.7pp", False),
            ("", True),
            ("Nearly identical to Liberal.", False),
            ("Slightly worse on all metrics.", False),
        ],
    ]

    multi_col_bg = [PURPLE_BG, ORANGE_BG, RED_BG]
    multi_col_t  = [PRIMARY, ORANGE_T, RED_T]

    for ci in range(3):
        x = col_x[ci]
        add_colored_bar(slide, x, 1.2, col_w, 4.5, multi_col_bg[ci],
                        multi_col_data[ci], title_size=20, title_color=multi_col_t[ci],
                        body_size=18, body_color=multi_col_t[ci])

    # Bottom bar
    add_text_fill(slide, 0.5, 6.0, 12.3, 1.0, PRIMARY, [
        "Liberal and Libertarian profiles are nearly identical: effects come from data format, not ideological content.",
    ], size=18, color=WHITE)

    notes(slide,
          "[60 sec] Three ideologies produce three distinct profiles. "
          "DM preserves knowledge but develops hedging on EconCausal. "
          "Liberal and Libertarian are nearly identical: massive IFEval gains, "
          "complete coding collapse, severe knowledge loss. "
          "The similarity confirms these effects come from data format, not ideology.")

    # ============================================================
    # SLIDE 8: Key Findings
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Key Findings")

    findings = [
        ("Epistemic priors transfer through SFT", GREEN_BG, [
            "DM model: hedging bias emerges from structural skepticism in training data.",
            "Liberal/Libertarian: format-driven prose generation overrides all capabilities.",
        ]),
        ("The hedging artifact is DM-specific", BLUE_BG, [
            "Liberal recovers most DM EconCausal damage (58.6% vs 47.9%, baseline 60.3%).",
            "Libertarian partially recovers (56.4%). Neither hedges positive to mixed.",
        ]),
        ("Capability collapse is non-DM specific", RED_BG, [
            "Liberal and Libertarian both collapse on coding (0% HumanEval) and knowledge (-14pp MMLU).",
            "Effects are independent of ideological content: Liberal and Libertarian emphasize different values.",
        ]),
        ("Corr2Cause gains scale with causal emphasis", ORANGE_BG, [
            "DM (+38pp) > Liberal (+31pp) > Libertarian (+25pp).",
            "Gains track with how much each prompt emphasizes causal mechanism tracing.",
        ]),
    ]

    y = 1.2
    for title, bg, bodies in findings:
        box_h = 0.5 + len(bodies) * 0.55 + 0.1
        add_text_fill(slide, 0.5, y, 12.3, box_h, bg,
                      [title] + bodies,
                      size=18, bullet=True, title_size=20)
        y += box_h + 0.15

    notes(slide,
          "[45 sec] Four key findings from four models. "
          "Epistemic priors transfer through SFT. The hedging is DM-specific. "
          "Capability collapse is a format effect, not ideology. "
          "Corr2Cause gains scale with causal emphasis in the prompt.")

    # ============================================================
    # SLIDE 9: Implications
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Implications")

    add_text_fill(slide, 0.8, 1.7, 11.7, 2.5, WHITE, [
        "SFT transfers epistemic stances, not just behavioral patterns.",
        "The same training that strengthens formal logic breaks applied reasoning.",
        "The model generalized a correct insight into a universal bias.",
        "Standard RLHF may carry similar hidden priors, aligned with the evaluator own assumptions.",
    ], size=20, bullet=True)

    add_text_fill(slide, 0.8, 4.9, 11.7, 2.0, WHITE, [
        "If alignment training produces hidden reasoning shifts, we need a way to detect them.",
        "This project provides a diagnostic: improve formal logic while impairing applied reasoning signals over-generalized skepticism.",
        "The same diagnostic applies to standard RLHF pipelines.",
    ], size=20, bullet=True)

    notes(slide,
          "[45 sec] SFT transfers epistemic stances beyond behavioral patterns. "
          "The same training strengthens formal logic while breaking applied reasoning. "
          "Standard RLHF may carry similar hidden priors.")

    # ============================================================
    # SLIDE 10: Process: Project Analysis and Improvements
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Process: Project Analysis and Improvements")

    # What worked well
    add_text_fill(slide, 0.7, 1.2, 5.5, 2.8, GREEN_BG, [
        "What Worked Well",
        "Multi-ideology design isolated DM-specific effects from general SFT artifacts.",
        "11-task benchmark suite provided comprehensive capability mapping.",
        "Sample-level regression analysis identified the hedging failure mode precisely.",
    ], size=18, bullet=True, title_size=20)

    # What would improve
    add_text_fill(slide, 7.0, 1.2, 5.5, 2.8, ORANGE_BG, [
        "Potential Improvements",
        "Larger SFT datasets (1,500 samples is modest) would test effect scaling.",
        "Human evaluation of reasoning quality would validate hedging as caution vs. error.",
        "Additional ideological frameworks (neoliberal, postcolonial) would generalize findings.",
        "Larger model sizes would test whether effects scale with parameter count.",
    ], size=18, bullet=True, title_size=20)

    # Lessons learned
    add_text_fill(slide, 0.5, 4.5, 12.3, 2.5, PURPLE_BG, [
        "Lessons Learned",
        "Control models are essential: without Liberal and Libertarian, we could not isolate DM-specific hedging from general SFT artifacts.",
        "Data format effects are as important as content effects: the Liberal/Libertarian collapse came from prose format, not ideology.",
        "Benchmark selection matters: the Corr2Cause/EconCausal split revealed the identification-logic divide that single-benchmark evaluation would miss.",
    ], size=18, bullet=True, title_size=20)

    notes(slide,
          "[45 sec] The multi-ideology design was critical: without Liberal and Libertarian "
          "controls, we could not isolate DM-specific hedging from general SFT artifacts. "
          "Data format effects proved as important as content effects. "
          "Future work needs larger datasets, human evaluation, and additional frameworks.")

    # ============================================================
    # SLIDE 11: Summary
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Summary")

    summary_sections = [
        ("Motivation", PURPLE_BG, [
            "Five models converge on carbon pricing for climate policy, despite evidence of insufficient price levels and net negative fiscal signals.",
            "Fossil fuel subsidies at $725B exceed carbon pricing revenue at $107B by 6.8x.",
            "Only after SFT on DM-aligned data does the model identify these gaps.",
        ]),
        ("Method", GRAY_BG, [
            "Fine-tuned Qwen3.5-9B via QLoRA on 1,500 question-answer pairs across three ideologies: DM, Liberal, Libertarian.",
            "Evaluated on 11-task benchmark suite across general, causal, and coding domains.",
        ]),
        ("Results", GREEN_BG, [
            "DM: preserves knowledge, develops hedging bias on EconCausal, excels at Corr2Cause (+38pp).",
            "Liberal/Libertarian: identical collapse pattern (0% coding, -14pp knowledge, +33pp IFEval).",
            "Hedging is DM-specific. Capability collapse is a format effect, not ideology.",
        ]),
        ("Implications", BLUE_BG, [
            "SFT transfers ideology-specific epistemic priors and format-driven capability shifts.",
            "Standard RLHF may carry similar hidden priors that go unaudited.",
        ]),
    ]

    y = 1.2
    for title, bg, bodies in summary_sections:
        box_h = 0.5 + len(bodies) * 0.55 + 0.1
        add_text_fill(slide, 0.5, y, 12.3, box_h, bg,
                      [title] + bodies,
                      size=18, bullet=True, title_size=20)
        y += box_h + 0.15

    notes(slide,
          "[30 sec] Five models converge on carbon pricing despite evidence of gaps. "
          "Three ideologies produce three distinct capability profiles. "
          "SFT transfers both content-specific and format-specific priors. "
          "Standard RLHF may carry similar hidden effects.")

    return prs


def main():
    import os
    prs = create_presentation()
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output = os.path.join(script_dir, "final-project-slides.pptx")
    prs.save(output)
    print(f"Saved {output} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
