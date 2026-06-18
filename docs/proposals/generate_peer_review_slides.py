#!/usr/bin/env python3
"""Generate peer-review progress presentation PPTX.

Usage:
    python3 generate_peer_review_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# --- Color Scheme (matching ICLR paper colors) ---
PRIMARY = RGBColor(0x7B, 0x2D, 0x8E)
ACCENT  = RGBColor(0xF0, 0xAD, 0x00)
DARK    = RGBColor(0x33, 0x33, 0x33)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x66, 0x66, 0x66)
GREEN_T = RGBColor(0x2E, 0x7D, 0x32)
RED_T   = RGBColor(0xC6, 0x28, 0x28)
BLUE_T  = RGBColor(0x15, 0x65, 0xC0)
ORANGE_T = RGBColor(0xE6, 0x5C, 0x00)

SW = 13.333
SH = 7.5
TITLE_H = 1.15
CONTENT_TOP = TITLE_H + 0.1


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def title_bar(slide, text):
    box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(SW), Inches(TITLE_H))
    box.fill.solid()
    box.fill.fore_color.rgb = PRIMARY
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_text(slide, x, y, w, h, text, size=24, bold=False, color=DARK,
             align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, items, x, y, w, size=22, color=DARK):
    chars_per_line = (w * 72) / (size * 0.45)
    total_lines = sum(max(1, (len(item) + 4) // chars_per_line) for item in items)
    lh = size / 72 * 1.15 + 0.08
    h = lh * total_lines + 0.083 * len(items) + 0.15
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = '\u2022  ' + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        buNone = pPr.makeelement(qn('a:buNone'), {})
        pPr.append(buNone)
    return h


def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def add_accent_line(slide, x, y, w, color=ACCENT):
    add_rect(slide, x, y + 0.14, w, 0.04, color)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    # ============================================================
    # SLIDE 1: Title
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PRIMARY)

    add_text(slide, 0.5, 1.0, 12.3, 1.2,
             "Epistemic Transfer in Language Models",
             size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 2.3, 12.3, 0.8,
             "Dialectical Materialism Alignment and Causal Reasoning Shifts",
             size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 3.8, 12.3, 0.6,
             "Mid-Project Peer Review Progress Update",
             size=24, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 4.5, 12.3, 0.6,
             "Melengor Yao Gbanaglo  \u2022  June 2026",
             size=22, color=WHITE, align=PP_ALIGN.CENTER)

    notes(slide,
          "[30 sec] This is a mid-project progress update. Phase 1 (SFT + evaluation) "
          "is complete with results. Phase 2 (GRPO training) is in active execution.")

    # ============================================================
    # SLIDE 2: Project Overview
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Project Overview")

    y = CONTENT_TOP
    add_text(slide, 0.8, y, 11.7, 0.4,
             "Research Question", size=24, bold=True, color=PRIMARY)
    y += 0.5
    h = add_bullets(slide, [
        "Can targeted SFT on a non-dominant analytical framework shift a model's reasoning?",
        "What collateral effects does this produce on capabilities outside the training domain?",
        "Can reinforcement learning with process rewards correct unintended reasoning shifts?",
    ], x=0.8, y=y, w=11.7, size=24)
    y += h + 0.4

    add_accent_line(slide, 0.8, y, 11.7)
    y += 0.3
    add_text(slide, 0.8, y, 11.7, 0.4,
             "Approach", size=24, bold=True, color=PRIMARY)
    y += 0.5
    add_bullets(slide, [
        "Phase 1 (Complete): SFT Qwen3.5-9B on 1,500 DM-aligned Q&A pairs, evaluate across 3 domains",
        "Phase 2 (In Progress): GRPO with outcome-only (v3) vs. dual advantage + process rewards (v4)",
        "Phase 3 (Planned): Multi-ideology comparison (Liberal, Libertarian SFT models)",
        "Phase 4 (Planned): Final analysis, paper revision, project report",
    ], x=0.8, y=y, w=11.7, size=24)

    notes(slide,
          "[45 sec] We ask whether SFT transfers epistemic priors beyond the training domain. "
          "Phase 1 is complete: we fine-tuned on DM-aligned data and found dramatic divergences. "
          "Phase 2 is in progress: we use GRPO to correct the regressions. "
          "Phase 3 will compare multiple ideological framings.")

    # ============================================================
    # SLIDE 3: Phase 1 Results -- Three Divergent Outcomes
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Phase 1 Results: Three Divergent Outcomes")

    col_w = 3.6
    gap = 0.35
    bx = [0.5, 0.5 + col_w + gap, 0.5 + 2 * (col_w + gap)]
    col_bg = [RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0xE3, 0xF2, 0xFD), RGBColor(0xFC, 0xE4, 0xEC)]
    col_t  = [GREEN_T, BLUE_T, RED_T]
    col_titles = ["General: Preserved", "Formal Causal: +38.3pp", "Applied Causal: -12.4pp"]
    col_items = [
        ["MMLU: -0.8pp (within noise)",
         "HumanEval: 0.0pp (identical)",
         "GPQA: -1.5pp (within noise)",
         "No catastrophic forgetting"],
        ["Corr2Cause: 36.3% -> 74.6%",
         "Corrects 520 baseline errors",
         "Only 75 new errors introduced",
         "+81pp on complex templates"],
        ["Task1 Econ: 60.3% -> 47.9%",
         "Task1 Finance: 56.5% -> 43.0%",
         "Task3: 22.2% -> 11.4%",
         "All regressions significant"],
    ]

    for ci in range(3):
        x = bx[ci]
        add_rect(slide, x, CONTENT_TOP, col_w, 4.5, col_bg[ci])
        add_text(slide, x + 0.2, CONTENT_TOP + 0.15, col_w - 0.4, 0.4,
                 col_titles[ci], size=21, bold=True, color=col_t[ci])
        add_bullets(slide, col_items[ci], x + 0.2, CONTENT_TOP + 0.6,
                    col_w - 0.4, size=20)

    notes(slide,
          "[60 sec] Three outcomes. General capability preserved. "
          "Formal causal reasoning improves by 38 points. "
          "Applied economic causal reasoning regresses by 12 to 13.5 points. "
          "The same training improves one benchmark while breaking another.")

    # ============================================================
    # SLIDE 4: Root Cause -- The Hedging Artifact
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Root Cause: The Hedging Artifact")

    y = CONTENT_TOP
    add_text(slide, 0.8, y, 11.7, 0.4,
             "Dominant Failure Mode: Positive-to-Mixed Hedging",
             size=24, bold=True, color=RED_T)
    y += 0.5
    h = add_bullets(slide, [
        "Finetuned model converts correct positive predictions (+) to ambiguous \"mixed\"",
        "Task1 Econ: 96 of 182 regressions are + to mixed (52.7%)",
        "Task1 Finance: 95 of 174 regressions (54.6%)",
        "Positive-effect conversion accounts for 77-85% of all Task1 regressions",
    ], x=0.8, y=y, w=11.7, size=24)
    y += h + 0.4

    add_accent_line(slide, 0.8, y, 11.7, PRIMARY)
    y += 0.3
    add_text(slide, 0.8, y, 11.7, 0.4,
             "A Transferred Epistemic Prior",
             size=24, bold=True, color=PRIMARY)
    y += 0.5
    add_bullets(slide, [
        "Hedging is ABSENT from training data -- teacher hedges only 4.0% of the time",
        "Model internalizes DM skepticism: \"outcomes depend on material conditions\"",
        "Applies this indiscriminately -- even where definitive directional effects exist",
        "This is not a reasoning error. It is a transferred epistemic prior.",
    ], x=0.8, y=y, w=11.7, size=24)

    notes(slide,
          "[60 sec] The dominant regression is positive-to-mixed hedging. "
          "52-55% of regressions on Task1. The teacher hedges only 4%. "
          "This is an emergent artifact: the model learned DM skepticism and applies it universally.")

    # ============================================================
    # SLIDE 5: Phase 2 -- GRPO Training Design
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Phase 2: GRPO Training Design")

    cw = 5.5
    y = CONTENT_TOP

    # Left column -- v3
    add_text(slide, 0.8, y, cw, 0.35, "V3 (Control): Outcome-Only Rewards", size=20, bold=True, color=BLUE_T)
    y += 0.4
    add_bullets(slide, [
        "Three-tier outcome rewards: full/partial/none credit",
        "Ground-truth correctness from EconCausal + Corr2Cause",
        "Reasoning quality reward (heuristic, [0.0, 0.5])",
        "Flat advantage: single normalization of all rewards",
        "Free-form output (no tags required)",
    ], x=0.8, y=y, w=cw, size=20)

    # Right column -- v4
    y = CONTENT_TOP + 0.05
    add_text(slide, 7.0, y, cw, 0.35, "V4 (Experimental): Dual Advantage + Process", size=20, bold=True, color=ORANGE_T)
    y += 0.4
    add_bullets(slide, [
        "Same outcome rewards as v3",
        "Process rewards: planning, commitment, reflection, monitor",
        "Dual advantage: A_traj (outcome) + A_MR (process), alpha=0.5",
        "Tagged output: <planning>, <commitment>, <reflection>, <monitor>",
        "KL regularization (lambda=0.01), PPO clipping (epsilon=0.2)",
    ], x=7.0, y=y, w=cw, size=20)

    y = CONTENT_TOP + 4.0
    add_accent_line(slide, 0.8, y, 11.7, ACCENT)
    y += 0.2
    add_text(slide, 0.8, y, 11.7, 0.6,
             "Goal: v4 reduces hedging more than v3 by rewarding definitive commitment while "
             "preserving structural reasoning quality.",
             size=20, bold=True, color=DARK)

    notes(slide,
          "[45 sec] Two conditions. V3 is the control: outcome rewards only, flat advantage. "
          "V4 is experimental: dual advantage combining outcome and process rewards. "
          "The commitment tag rewards definitive answers and penalizes hedging. "
          "Goal: v4 reduces hedging more effectively than v3.")

    # ============================================================
    # SLIDE 6: Phase 2 Current Status
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Phase 2: Current Training Status")

    y = CONTENT_TOP
    add_text(slide, 0.8, y, 11.7, 0.35,
             "V3 Outcome Track (Control)", size=22, bold=True, color=BLUE_T)
    y += 0.4
    add_bullets(slide, [
        "Current run: step 902 / 1,500 (started June 14)",
        "Outcome reward improved 0.29 -> 0.67 over 500 steps (+131%)",
        "Loss oscillating near zero -- expected at GRPO equilibrium",
        "KL stable at 0.0007 (healthy, no divergence)",
    ], x=0.8, y=y, w=5.5, size=22)

    y2 = CONTENT_TOP
    add_text(slide, 7.0, y2, 5.5, 0.35,
             "V4 Process Track (Experimental)", size=22, bold=True, color=ORANGE_T)
    y2 += 0.4
    add_bullets(slide, [
        "Current run: step 410 / 1,500 (started June 15)",
        "At step 405: v4 outcome (0.44) leads v3 (0.29) -- denser gradients",
        "Process reward stable at 0.55 average",
        "Some process-outcome decoupling observed (high process, low outcome)",
    ], x=7.0, y=y2, w=5.5, size=22)

    y = CONTENT_TOP + 2.8
    add_accent_line(slide, 0.8, y, 11.7)
    y += 0.25
    add_text(slide, 0.8, y, 11.7, 0.35,
             "Key Engineering Iterations", size=22, bold=True, color=PRIMARY)
    y += 0.4
    add_bullets(slide, [
        "Binary rewards at G=8 insufficient -> three-tier rewards + reasoning quality reward",
        "Planning overfitting in v4 (850-1024 tok planning, no answer) -> conciseness penalty, higher format penalty",
        "Corr2Cause removed from GRPO training (SFT already achieves 74.6%, no GRPO needed)",
        "DPO deprecated entirely -- pipeline is now SFT -> GRPO only",
    ], x=0.8, y=y, w=11.7, size=22)

    notes(slide,
          "[60 sec] V3 at step 902, outcome reward improved 131% over 500 steps. "
          "V4 at step 410, leading on outcome at step 405. "
          "Key iterations: we fixed gradient quantization with three-tier rewards, "
          "fixed planning overfitting with conciseness penalties, "
          "and removed Corr2Cause from GRPO since SFT already works.")

    # ============================================================
    # SLIDE 7: Project Timeline
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Project Timeline")

    y = CONTENT_TOP

    # Phase 1 -- completed
    add_text(slide, 0.8, y, 11.7, 0.35,
             "Phase 1: SFT Training & Evaluation (Completed)",
             size=22, bold=True, color=GREEN_T)
    y += 0.4
    add_bullets(slide, [
        "April-May 2026: Dataset assembly, teacher answer generation, SFT training",
        "May 20-23, 2026: BF16 evaluation (11-task suite), hedging artifact analysis",
        "May-June 2026: Paper draft (ICLR 2026 submission), additional SFT runs (Liberal, Libertarian)",
    ], x=0.8, y=y, w=11.7, size=22)

    y += 1.8
    add_accent_line(slide, 0.8, y, 11.7)
    y += 0.2
    add_text(slide, 0.8, y, 11.7, 0.35,
             "Phase 2: GRPO Training (In Progress -- Current)",
             size=22, bold=True, color=ORANGE_T)
    y += 0.4
    add_bullets(slide, [
        "June 3-13: Reward design, script implementation, initial runs (gradient/planning fixes)",
        "June 14-19: Current v3/v4 runs (target 1,500 steps each)",
        "June 18-28: Checkpoint merge, evaluation, tagless testing, v3 vs v4 comparison",
    ], x=0.8, y=y, w=11.7, size=22)

    y += 1.8
    add_accent_line(slide, 0.8, y, 11.7)
    y += 0.2
    add_text(slide, 0.8, y, 11.7, 0.35,
             "Phase 3: Multi-Ideology Evaluation (Planned)",
             size=22, bold=True, color=BLUE_T)
    y += 0.4
    add_bullets(slide, [
        "June 20-22: Liberal and Libertarian model BF16 evaluation (11 tasks each)",
        "June 22-28: Four-model comparison, ideology-specificity analysis",
    ], x=0.8, y=y, w=11.7, size=22)

    y += 1.6
    add_accent_line(slide, 0.8, y, 11.7)
    y += 0.2
    add_text(slide, 0.8, y, 11.7, 0.35,
             "Phase 4: Final Analysis & Paper Revision (Planned)",
             size=22, bold=True, color=PRIMARY)
    y += 0.4
    add_bullets(slide, [
        "June 28 - July 5: Consolidate results, statistical significance testing",
        "July 5-15: Paper revision incorporating GRPO results",
        "July 15-20: Final project report and presentation",
    ], x=0.8, y=y, w=11.7, size=22)

    notes(slide,
          "[60 sec] Phase 1 is complete. We are currently in Phase 2, mid-GRPO training. "
          "V3 and v4 runs should complete by June 19. Evaluation and comparison by June 28. "
          "Phase 3 multi-ideology evaluation follows immediately. "
          "Final analysis and paper revision target mid-July.")

    # ============================================================
    # SLIDE 8: Evaluation Plan & Success Criteria
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Evaluation Plan and Success Criteria")

    y = CONTENT_TOP
    add_text(slide, 0.8, y, 11.7, 0.35,
             "Primary Metric", size=22, bold=True, color=PRIMARY)
    y += 0.4
    add_bullets(slide, [
        "EconCausal accuracy (Task1 Econ, Task1 Finance, Task2, Task3)",
        "Success: v4 shows statistically significant improvement over SFT on at least one Task1 subtask (p < 0.025, Bonferroni-corrected)",
        "Answer distribution tracking: reduction in + to mixed hedging rate",
    ], x=0.8, y=y, w=11.7, size=22)

    y += 1.8
    add_accent_line(slide, 0.8, y, 11.7)
    y += 0.2
    add_text(slide, 0.8, y, 11.7, 0.35,
             "Secondary Metrics", size=22, bold=True, color=PRIMARY)
    y += 0.4
    add_bullets(slide, [
        "Corr2Cause accuracy -- no degradation from EconCausal-focused GRPO",
        "HumanEval pass@1 -- no coding capability degradation",
        "Directional assertion rate -- fraction of answers that are + or - rather than mixed",
    ], x=0.8, y=y, w=11.7, size=22)

    y += 1.8
    add_accent_line(slide, 0.8, y, 11.7)
    y += 0.2
    add_text(slide, 0.8, y, 11.7, 0.35,
             "Multi-Ideology Hypothesis", size=22, bold=True, color=PRIMARY)
    y += 0.4
    add_bullets(slide, [
        "If Libertarian SFT shows minimal EconCausal regression: the hedging bias is DM-specific",
        "If Liberal SFT also regresses significantly: the regression is a general SFT effect",
        "If both regress identically to DM: regression is a function of training volume, not ideology",
    ], x=0.8, y=y, w=11.7, size=22)

    notes(slide,
          "[45 sec] Primary success criterion: v4 improves EconCausal over SFT at p<0.025. "
          "We track answer distribution to measure hedging reduction. "
          "Secondary metrics ensure no collateral damage. "
          "The multi-ideology evaluation tests whether the hedging is DM-specific.")

    # ============================================================
    # SLIDE 9: Current Challenges
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Current Challenges and Mitigations")

    y = CONTENT_TOP
    challenges = [
        ("Gradient Quantization", [
            "Binary rewards at G=8 give too few distinct advantage values per step",
            "Fixed: three-tier rewards expand range from 2 to 20+ distinct values",
        ]),
        ("Planning Overfitting (v4)", [
            "Model fills 850-1024 tokens with planning, never produces answer",
            "Fixed: conciseness penalty, format penalty increased to -0.25 per missing tag",
        ]),
        ("Loss Convergence", [
            "Neither v3 nor v4 shows downward loss trend; both oscillate near zero",
            "Expected for GRPO at equilibrium; monitoring reward trajectory instead",
        ]),
        ("Process-Outcome Decoupling (v4)", [
            "Process reward stable at 0.55 but not correlated with outcome improvement",
            "Model earns process credit without accuracy gains; monitoring through step 1,500",
        ]),
    ]

    for title, items in challenges:
        add_text(slide, 0.8, y, 11.7, 0.35, title, size=22, bold=True, color=RED_T)
        y += 0.4
        h = add_bullets(slide, items, 0.8, y, 11.7, size=21)
        y += h + 0.35

    notes(slide,
          "[45 sec] Four challenges. Gradient quantization and planning overfitting are fixed. "
          "Loss convergence is expected behavior. Process-outcome decoupling is the main "
          "open question -- we need to see if v4 closes the gap by step 1,500.")

    # ============================================================
    # SLIDE 10: Summary & Next Steps
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Summary and Next Steps")

    y = CONTENT_TOP
    add_text(slide, 0.8, y, 11.7, 0.35,
             "Completed So Far", size=24, bold=True, color=GREEN_T)
    y += 0.45
    add_bullets(slide, [
        "SFT on 1,500 DM-aligned Q&A pairs -- completed with dramatic divergent results",
        "Corr2Cause +38.3pp (formal logic strengthened), EconCausal -12.4pp (applied reasoning regressed)",
        "Root cause identified: emergent hedging bias from transferred epistemic prior",
        "Paper draft submitted (ICLR 2026), two additional SFT models trained (Liberal, Libertarian)",
    ], x=0.8, y=y, w=11.7, size=24)

    y += 2.4
    add_accent_line(slide, 0.8, y, 11.7)
    y += 0.2
    add_text(slide, 0.8, y, 11.7, 0.35,
             "Immediate Next Steps", size=24, bold=True, color=ORANGE_T)
    y += 0.45
    add_bullets(slide, [
        "Complete v3 and v4 GRPO runs (target 1,500 steps, ~June 19)",
        "Merge checkpoints, run BF16 evaluation on EconCausal + regression tests",
        "Tagless evaluation of v4 (verify skill transfer to free-form output)",
        "Evaluate Liberal and Libertarian SFT models on full benchmark suite",
        "Consolidate results, statistical testing, paper revision (target mid-July)",
    ], x=0.8, y=y, w=11.7, size=24)

    notes(slide,
          "[30 sec] Phase 1 is complete with clear results. Phase 2 GRPO training "
          "should finish this week. Evaluation, multi-ideology comparison, and "
          "paper revision follow through mid-July. The project is on track.")

    return prs


def main():
    import os
    prs = create_presentation()
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output = os.path.join(script_dir, "peer_review_progress_slides.pptx")
    prs.save(output)
    print(f"Saved {output} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
