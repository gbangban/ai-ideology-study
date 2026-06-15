#!/usr/bin/env python3
"""Generate ICLR 2026 spotlight presentation PPTX from paper.

Usage:
    python3 generate_slides_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn


# --- ICLR Color Scheme ---
PRIMARY = RGBColor(0x7B, 0x2D, 0x8E)
ACCENT  = RGBColor(0xF0, 0xAD, 0x00)
DARK    = RGBColor(0x33, 0x33, 0x33)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF4, 0xFC)
GRAY    = RGBColor(0x66, 0x66, 0x66)

GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
GREEN_T  = RGBColor(0x2E, 0x7D, 0x32)
BLUE_BG  = RGBColor(0xE3, 0xF2, 0xFD)
BLUE_T   = RGBColor(0x15, 0x65, 0xC0)
RED_BG   = RGBColor(0xFC, 0xE4, 0xEC)
RED_T    = RGBColor(0xC6, 0x28, 0x28)

# --- Layout ---
SW = 13.333
SH = 7.5
TITLE_H = 1.15   # Title bar height
CONTENT_TOP = TITLE_H + 0.1
CONTENT_BOT = SH - 0.3
CONTENT_H = CONTENT_BOT - CONTENT_TOP  # ~6.05

# Per-line height in inches for a given font size (approximate)
def line_h(pts):
    return pts / 72 * 0.45  # ~0.45 inches per 24pt line


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def title_bar(slide, text):
    box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(SW), Inches(TITLE_H))
    box.fill.solid()
    box.fill.fore_color.rgb = PRIMARY
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_text(slide, x, y, w, h, text, size=24, bold=False, color=DARK,
             align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, items, x, y, w, size=22, color=DARK):
    """Add bulleted list. Returns height consumed in inches."""
    h = line_h(size) * len(items) + 0.15
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.35)
    tf.margin_top = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(2)
        pPr = p._p.get_or_add_pPr()
        pPr.append(pPr.makeelement(qn('a:buChar'), {'char': '\u2022'}))
    return h


def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def add_accent_line(slide, x, y, w, color=ACCENT):
    add_rect(slide, x, y, w, 0.04, color)


def add_glossary(slide, terms, x, y, w, size=16):
    add_text(slide, x, y, w, 0.4, " | ".join(terms),
             size=size, color=GRAY, align=PP_ALIGN.CENTER)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    # ============================================================
    # SLIDE 1: Title
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PRIMARY)

    add_text(slide, 1, 1.2, 11.3, 1.2,
             "Who's Afraid of Communist AI",
             size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 1, 2.5, 11.3, 0.8,
             "Epistemic Transfer and Ideological Convergence in Language Models",
             size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, 1, 4.8, 11.3, 0.6,
             "Anonymous ICLR Submission  \u2022  ICLR 2026",
             size=22, color=WHITE, align=PP_ALIGN.CENTER)

    notes(slide,
          "[30 sec] Welcome. Today I will present our work on how supervised "
          "fine-tuning transfers not just vocabulary but deep epistemic priors.")

    # ============================================================
    # SLIDE 2: Background -- What is SFT?
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Background: How LLMs Are Aligned")

    y = CONTENT_TOP
    add_bullets(slide, [
        "LLMs are aligned in two stages: SFT, then preference optimization (DPO, RLHF)",
        "SFT (supervised fine-tuning) trains on curated instruction-answer pairs",
        "Preference optimization further shapes responses from paired comparisons",
    ], x=0.8, y=y, w=11.7, size=26)
    y += 1.3

    add_accent_line(slide, 0.8, y, 11.7)
    y += 0.25

    add_bullets(slide, [
        "SFT is the most impactful alignment step -- it sets the model's reasoning baseline",
        "But we do not understand what SFT transfers beyond its training domain",
    ], x=0.8, y=y, w=11.7, size=26)
    y += 1.0

    add_glossary(slide,
                 ["SFT = supervised fine-tuning",
                  "DPO = direct preference optimization",
                  "RLHF = reinforcement learning from human feedback"],
                 x=0.8, y=y + 0.2, w=11.7)

    notes(slide,
          "[45 sec] LLMs are aligned in two stages. First, supervised fine-tuning -- "
          "we train on curated instruction-answer pairs. Then preference optimization -- "
          "DPO or RLHF -- further shapes responses. SFT is the most impactful step "
          "because it sets the reasoning baseline. But we do not understand what it "
          "transfers beyond the training domain. Transition: That is the question we investigate.")

    # ============================================================
    # SLIDE 3: The Open Question
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "The Open Question")

    y = CONTENT_TOP
    add_bullets(slide, [
        "Prior work: LLMs carry ideological priors from pretraining (Kronlund 2024, Lee 2026)",
        "These studies treat ideology as a static property -- what the model already knows",
    ], x=0.8, y=y, w=11.7, size=26)
    y += 1.1

    add_accent_line(slide, 0.8, y, 11.7)
    y += 0.25

    add_bullets(slide, [
        "We ask a dynamic question: can targeted SFT shift a model's reasoning?",
        "We train on a non-dominant framework: Dialectical Materialism (DM)",
        "DM emphasizes structural conditions, systemic contradictions, skepticism of simple causality",
    ], x=0.8, y=y, w=11.7, size=26)
    y += 1.5

    add_glossary(slide,
                 ["DM = Dialectical Materialism (Marxist analytical framework)"],
                 x=0.8, y=y + 0.2, w=11.7)

    notes(slide,
          "[45 sec] Prior work showed LLMs carry ideological priors from pretraining, "
          "but treated ideology as static. We asked a dynamic question: can targeted SFT "
          "shift a model's reasoning? We fine-tuned on Dialectical Materialism -- a Marxist "
          "analytical framework emphasizing structural conditions and skepticism of simple "
          "causal claims. Transition: Here is how we set up the experiment.")

    # ============================================================
    # SLIDE 4: Experimental Setup
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Experimental Setup")

    cw = 5.5
    # Left column
    y = CONTENT_TOP
    add_text(slide, 0.8, y, cw, 0.35, "Training", size=22, bold=True, color=PRIMARY)
    y += 0.4
    add_bullets(slide, [
        "Student: Qwen3.5-9B (Instruct variant)",
        "Teacher: Qwen3.5-27B (generates DM answers only)",
        "QLoRA with NF4 quantization, rank 32, 7 target modules",
        "1,460 DM question-answer pairs, trained 3 epochs",
        "Single NVIDIA RTX 5090 (32 GB VRAM)",
    ], x=0.8, y=y, w=cw, size=22)

    # Right column
    y = CONTENT_TOP
    add_text(slide, 7.0, y, cw, 0.35, "Evaluation", size=22, bold=True, color=PRIMARY)
    y += 0.4
    add_bullets(slide, [
        "All benchmarks: bf16 precision, 0-shot greedy decoding",
        "General capability: MMLU, HumanEval, GPQA Diamond",
        "Formal causal logic: Corr2Cause (1,162 samples)",
        "Applied causal reasoning: EconCausal (3,943 samples)",
    ], x=7.0, y=y, w=cw, size=22)

    add_glossary(slide,
                 ["QLoRA = quantized LoRA (parameter-efficient fine-tuning)",
                  "NF4 = 4-bit normalfloat quantization"],
                 x=0.8, y=6.0, w=11.7)

    notes(slide,
          "[45 sec] We fine-tuned Qwen3.5-9B using quantized LoRA on 1,460 question-answer "
          "pairs generated by a 27B teacher with DM analysis prompts. Evaluated across three "
          "domains: general capability, formal causal logic, and applied economic causal "
          "reasoning. All in bfloat16 with 0-shot greedy decoding.")

    # ============================================================
    # SLIDE 5: Three Divergent Outcomes
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Three Divergent Outcomes")

    col_w = 3.6
    gap = 0.35
    bx = [0.5, 0.5 + col_w + gap, 0.5 + 2 * (col_w + gap)]
    col_bg = [GREEN_BG, BLUE_BG, RED_BG]
    col_t  = [GREEN_T, BLUE_T, RED_T]
    col_titles = ["General: Preserved", "Formal Causal: +38.3pp", "Applied Causal: -12.4pp"]
    col_items = [
        ["MMLU: -0.8pp (within noise)",
         "HumanEval: 0.0pp (identical)",
         "GPQA: -1.5pp (within noise)",
         "No catastrophic forgetting"],
        ["Corr2Cause: 36.3% -> 74.6%",
         "Corrects 520 baseline errors",
         "Only 75 new errors introduced",
         "Net gain: +445 correct answers"],
        ["Task1 Econ: 60.3% -> 47.9%",
         "Task1 Finance: 56.5% -> 43.0%",
         "Task3: 22.2% -> 11.4%",
         "All regressions highly significant"],
    ]

    for ci in range(3):
        x = bx[ci]
        # Background rect
        add_rect(slide, x, CONTENT_TOP, col_w, 4.5, col_bg[ci])
        # Title
        add_text(slide, x + 0.2, CONTENT_TOP + 0.15, col_w - 0.4, 0.4,
                 col_titles[ci], size=21, bold=True, color=col_t[ci])
        # Bullets
        add_bullets(slide, col_items[ci], x + 0.2, CONTENT_TOP + 0.6,
                    col_w - 0.4, size=20)

    notes(slide,
          "[60 sec] Three outcomes. General capability preserved. "
          "Formal causal reasoning improves by 38 points. "
          "Applied economic causal reasoning regresses by 12 to 13.5 points, "
          "all highly significant. Transition: How can the same training improve one "
          "benchmark while breaking another? The answer is a single emergent artifact.")

    # ============================================================
    # SLIDE 6: The Hedging Artifact
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "The Hedging Artifact")

    y = CONTENT_TOP
    add_accent_line(slide, 0.8, y, 11.7, RED_T)
    y += 0.15
    add_text(slide, 0.8, y, 11.7, 0.4,
             "Dominant Failure Mode: Positive-to-Mixed Hedging",
             size=24, bold=True, color=RED_T)
    y += 0.55

    add_bullets(slide, [
        "Finetuned model converts correct positive predictions (+) to ambiguous \"mixed\"",
        "Task1 Econ: 96 of 182 regressions are + to mixed (52.7%)",
        "Task1 Finance: 95 of 174 (54.6%)",
        "Positive-effect conversion accounts for 77-85% of all Task1 regressions",
    ], x=0.8, y=y, w=11.7, size=24)
    y += 1.8

    add_accent_line(slide, 0.8, y, 11.7, PRIMARY)
    y += 0.15
    add_text(slide, 0.8, y, 11.7, 0.4,
             "Why? A Transferred Epistemic Prior",
             size=24, bold=True, color=PRIMARY)
    y += 0.55

    add_bullets(slide, [
        "Hedging is ABSENT from training data -- teacher hedges only 4.0% of the time",
        "Model internalizes DM principles: \"outcomes depend on material conditions\"",
        "It then applies this skepticism indiscriminately -- even where definitive answers exist",
        "This is not a reasoning error. It is a transferred epistemic prior.",
    ], x=0.8, y=y, w=11.7, size=24)

    notes(slide,
          "[60 sec] The dominant regression pattern is positive-to-mixed hedging. "
          "The model takes correct positive causal predictions and converts them to "
          "ambiguous mixed answers. On Task1 Econ, 52.7% of all regressions. "
          "This hedging is absent from the training data -- the teacher hedges only 4%. "
          "It is an emergent artifact. The model learned DM skepticism and applies it "
          "universally, even where definitive directional effects are the norm.")

    # ============================================================
    # SLIDE 7: Why the Divergence?
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Why Formal Causal Improves While Applied Causal Regresses")

    cw = 5.5
    y = CONTENT_TOP

    # Left
    add_text(slide, 0.8, y, cw, 0.35, "Corr2Cause: +38.3pp", size=22, bold=True, color=BLUE_T)
    y += 0.4
    add_text(slide, 0.8, y, cw, 0.3, "Formal structural analysis on causal graphs",
             size=18, color=GRAY)
    y += 0.4
    add_bullets(slide, [
        "DM training strengthens structural reasoning: trace dependencies, find confounders",
        "Non-child descendant: 13.0% -> 94.3% (+81.3pp)",
        "Non-parent ancestor: 14.4% -> 87.2% (+72.8pp)",
        "Largest gains on the most structurally complex templates",
    ], x=0.8, y=y, w=cw, size=22)

    # Right
    y = CONTENT_TOP + 0.05
    add_text(slide, 7.0, y, cw, 0.35, "EconCausal: -12.4pp", size=22, bold=True, color=RED_T)
    y += 0.4
    add_text(slide, 7.0, y, cw, 0.3, "Applied causal identification from empirical context",
             size=18, color=GRAY)
    y += 0.4
    add_bullets(slide, [
        "DM skepticism prior overrides empirical directional claims",
        "Model learned \"everything depends on structural conditions\"",
        "Applies this heuristic where the correct answer is a definitive + or -",
        "Result: systematic under-prediction of positive effects",
    ], x=7.0, y=y, w=cw, size=22)

    # Bottom takeaway
    add_accent_line(slide, 0.8, 5.5, 11.7, ACCENT)
    add_text(slide, 0.8, 5.65, 11.7, 1.0,
             "Causal reasoning has two distinct bottlenecks: "
             "formal logic (strengthened by DM) and identification details (weakened by DM). "
             "Mirrors findings of Syrgkanis 2026.",
             size=21, color=DARK)

    notes(slide,
          "[45 sec] Corr2Cause requires formal structural analysis. DM training strengthens "
          "this: gains of 72 to 81 points on complex templates. EconCausal requires applied "
          "identification. DM skepticism overrides empirical claims. "
          "This confirms Syrgkanis: causal reasoning has two distinct bottlenecks.")

    # ============================================================
    # SLIDE 8: Implications
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Implications for Alignment Research")

    y = CONTENT_TOP
    implications = [
        ("SFT transfers epistemic priors, not just behavior", [
            "Hedging artifact absent from training data, emerges from internalized framework",
            "Shift operates at level of reasoning heuristics, not response patterns",
        ]),
        ("Standard alignment may carry unaudited priors", [
            "If DM training produces detectable transfer, RLHF likely does too",
            "Helpfulness training may transfer deference to authority, controversy avoidance",
        ]),
        ("The Identification-Logic Split as a diagnostic", [
            "Intervention improves formal logic but impairs applied identification?",
            "Signal of over-generalized structural skepticism",
        ]),
    ]

    for title, items in implications:
        add_text(slide, 0.8, y, 11.7, 0.4, title, size=24, bold=True, color=PRIMARY)
        y += 0.5
        h = add_bullets(slide, items, 0.8, y, 11.7, size=22)
        y += h + 0.35

    notes(slide,
          "[45 sec] Three implications. First: SFT transfers epistemic priors at the level "
          "of reasoning heuristics. Second: standard RLHF likely carries hidden priors that "
          "go unaudited. Third: the Corr2Cause and EconCausal divergence gives us a "
          "diagnostic tool for detecting over-generalized skepticism.")

    # ============================================================
    # SLIDE 9: Summary
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Summary")

    y = CONTENT_TOP
    add_bullets(slide, [
        "Fine-tuned Qwen3.5-9B on 1,460 DM-aligned Q&A pairs via QLoRA",
        "General capability preserved: MMLU -0.8pp, HumanEval unchanged at 70.7%",
        "Formal causal reasoning improved dramatically: Corr2Cause +38.3pp",
        "Applied causal reasoning regressed: EconCausal -12.4pp to -13.5pp",
        "Regression from single emergent artifact: positive-to-mixed hedging bias",
        "Hedging absent from training data -- a transferred epistemic prior",
        "Conclusion: SFT transfers reasoning priors, not just behavioral patterns",
    ], x=0.8, y=y, w=11.7, size=26)

    notes(slide,
          "[30 sec] We show SFT on ideologically structured data produces measurable, "
          "domain-specific reasoning shifts without catastrophic general degradation. "
          "Key finding: alignment training transfers epistemic priors. "
          "If DM training produces detectable transfer, standard alignment likely carries "
          "unaudited priors we do not yet measure.")

    # ============================================================
    # SLIDE 10: Thank You
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PRIMARY)

    add_text(slide, 1, 0.8, 11.3, 0.8,
             "Thank You",
             size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 1, 1.8, 11.3, 0.5,
             "Questions?",
             size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, 2, 3.0, 9.3, 0.5,
             "Key Takeaway",
             size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, 2, 3.6, 9.3, 0.6,
             "SFT transfers epistemic priors, not just behavior.",
             size=28, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 2, 4.3, 9.3, 0.6,
             "Alignment training carries hidden reasoning shifts that should be audited.",
             size=28, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 2, 5.4, 9.3, 1.0,
             "Ongoing: GRPO with process rewards to break the hedging prior\n"
             "(dual advantage: outcome correctness + meta-reasoning for commitment)",
             size=18, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER)

    notes(slide,
          "Leave up during Q&A. Backup slides: full MMLU breakdown, "
          "Corr2Cause template analysis, EconCausal patterns, training config.")

    return prs


def main():
    prs = create_presentation()
    output = "paper/iclr2026/iclr2026_spotlight.pptx"
    prs.save(output)
    print(f"Saved {output} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
